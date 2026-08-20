"""
Integration and End-to-End API Tests for Fortified Enterprise Fleet API (v0.2.0).
Tests fail-closed JWT auth, sanitized error messages, dev-token endpoint,
checkpoint persistence, 3-way digest verification, approval_request_id binding, and audit trail.
"""
import base64
import hashlib
import json
from pathlib import Path
from typing import Any, Dict, List, Tuple
from uuid import uuid4
import pytest
from fastapi.testclient import TestClient
from fleet_api.main import app
from fleet_api.security import create_access_token
from fleet_governance_core.ports.intake_port import IntakePort

import fleet_api.deps as deps
from fleet_adapter_gcp import (
    InMemoryApprovalStore,
    InMemoryArtifactStorageAdapter,
    InMemoryAuditLog,
    InMemoryCheckpointStore,
    InMemoryMemoryStore,
    ThreadSafeDocumentResolver,
)
from fleet_adapter_local import SQLiteResumeContextStore

FIXTURES_DIR = Path(__file__).resolve().parent.parent.parent.parent / "fixtures"

from fleet_api.routers import dossiers

@pytest.fixture(autouse=True)
def reset_in_memory_stores():
    deps.approval_store = InMemoryApprovalStore()
    deps.audit_log = InMemoryAuditLog()
    deps.checkpoint_store = InMemoryCheckpointStore()
    deps.memory_store = InMemoryMemoryStore()
    deps.storage_adapter = InMemoryArtifactStorageAdapter()
    deps.document_resolver = ThreadSafeDocumentResolver()
    deps.resume_context_store = SQLiteResumeContextStore(":memory:")
    dossiers.CASES_DB.clear()
    yield

@pytest.fixture
def client():
    return TestClient(app)

@pytest.fixture
def auth_headers():
    token = create_access_token(
        tenant_id="tenant-acme-corp",
        sub="usr-safety-officer-01",
        roles=["safety_assessor", "approver"],
        email="safety@acme.com",
    )
    return {"Authorization": f"Bearer {token}"}

@pytest.fixture
def viewer_headers():
    token = create_access_token(
        tenant_id="tenant-acme-corp",
        sub="usr-viewer-01",
        roles=["viewer"],
    )
    return {"Authorization": f"Bearer {token}"}

def test_health_endpoint(client):
    resp = client.get("/v1/health")
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "healthy"
    assert data["version"] == "0.3.2"
    assert data["runtime_mode"] == "local_memory_emulation"
    assert "adapters" in data
    assert data["adapters"]["orchestrator"]["configured_mode"] in ("fake", "live")
    assert data["adapters"]["intake"]["configured_mode"] in ("fake", "live")

def test_ready_endpoint(client):
    resp = client.get("/v1/ready")
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "ready"
    assert data["version"] == "0.3.2"
    assert "adapters" in data
    assert data["adapters"]["intake"]["status"] == "ready"
    assert data["adapters"]["orchestrator"]["status"] == "ready"

def test_intake_error_handlers(client):
    from fleet_adapter_prodocux import IntakeConnectionError, IntakeTimeoutError
    from fleet_api.main import app
    from fastapi.testclient import TestClient

    # Test custom test routes to trigger exception handlers directly
    @app.get("/v1/test-timeout")
    def trigger_timeout():
        raise IntakeTimeoutError("upstream timeout")

    @app.get("/v1/test-unavailable")
    def trigger_unavailable():
        raise IntakeConnectionError("upstream connection error")

    c = TestClient(app)
    r_timeout = c.get("/v1/test-timeout")
    assert r_timeout.status_code == 504
    assert "timed out" in r_timeout.json()["detail"].lower()

    r_unavailable = c.get("/v1/test-unavailable")
    assert r_unavailable.status_code == 502
    assert "unavailable" in r_unavailable.json()["detail"].lower()



def test_dev_token_endpoint(client):
    resp = client.post("/v1/auth/dev-token", json={
        "tenant_id": "tenant-acme-corp",
        "sub": "usr-cso-steven-wu",
        "roles": ["safety_assessor", "approver", "cso"]
    })
    assert resp.status_code == 200
    data = resp.json()
    assert "access_token" in data
    assert data["token_type"] == "bearer"
    assert data["tenant_id"] == "tenant-acme-corp"

def test_missing_auth_token_returns_401(client):
    resp = client.post("/v1/dossiers/create", json={})
    assert resp.status_code == 401
    assert "Authentication credentials were not provided" in resp.json()["detail"]

def test_unauthorized_role_approval_returns_403_sanitized(client, viewer_headers):
    decision_payload = {
        "checkpoint_id": "chk-any-001",
        "run_id": "run-001",
        "approval_request_id": str(uuid4()),
        "idempotency_key": "idemp-test-01",
        "decision": "approved",
        "case_digest": "0" * 64,
        "plan_digest": "1" * 64,
    }
    resp = client.post("/v1/approval/decide", json=decision_payload, headers=viewer_headers)
    assert resp.status_code == 403
    assert "Forbidden: Insufficient privileges" in resp.json()["detail"]

def test_unknown_checkpoint_returns_404(client, auth_headers):
    decision_payload = {
        "checkpoint_id": "chk-non-existent-999",
        "run_id": "run-001",
        "approval_request_id": str(uuid4()),
        "idempotency_key": "idemp-test-01",
        "decision": "approved",
        "case_digest": "0" * 64,
        "plan_digest": "1" * 64,
    }
    resp = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp.status_code == 404
    assert "not found under tenant" in resp.json()["detail"]

def test_document_registration_endpoint(client, auth_headers):
    # Register document bytes
    raw_pdf = b"%PDF-1.4 Minimal Test Doc"
    import base64
    payload = {
        "doc_id": "doc-sds-test-01",
        "content_b64": base64.b64encode(raw_pdf).decode("utf-8"),
        "filename": "sds-test.pdf",
    }
    resp = client.post("/v1/dossiers/documents/register", json=payload, headers=auth_headers)
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "registered"
    assert data["doc_id"] == "doc-sds-test-01"
    assert data["size_bytes"] == len(raw_pdf)

def test_document_registration_security_validations(client, auth_headers):
    # 1. Path traversal in doc_id rejected
    resp = client.post(
        "/v1/dossiers/documents/register",
        json={"doc_id": "../evil_doc", "content_b64": "AAAA", "filename": "evil.pdf"},
        headers=auth_headers,
    )
    assert resp.status_code == 422

    # 2. Path traversal in filename rejected
    resp = client.post(
        "/v1/dossiers/documents/register",
        json={"doc_id": "valid-id", "content_b64": "AAAA", "filename": "../../evil.pdf"},
        headers=auth_headers,
    )
    assert resp.status_code == 422

    # 3. Invalid base64 encoding returns sanitized 400
    resp = client.post(
        "/v1/dossiers/documents/register",
        json={"doc_id": "valid-id", "content_b64": "NOT_VALID_BASE64!@#$", "filename": "test.pdf"},
        headers=auth_headers,
    )
    assert resp.status_code == 400
    assert resp.json()["detail"] == "Invalid Base64 payload encoding."

    # 4. Unsupported document format returns sanitized 400
    resp = client.post(
        "/v1/dossiers/documents/register",
        json={"doc_id": "valid-id", "content_b64": base64.b64encode(b"foo").decode("utf-8"), "filename": "malware.exe"},
        headers=auth_headers,
    )
    assert resp.status_code == 400
    assert resp.json()["detail"] == "Unsupported document format."

def test_full_dossier_lifecycle_happy_path(client, auth_headers):
    # 1. Create Dossier Case
    happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    sample_pdf_bytes = b"%PDF-1.4 Happy Path Document"
    sample_pdf_sha256 = hashlib.sha256(sample_pdf_bytes).hexdigest()
    happy_raw["supplier_documents"][0]["sha256"] = sample_pdf_sha256

    # Register document in tenant resolver
    reg_payload = {
        "doc_id": "doc-sds-001",
        "content_b64": base64.b64encode(sample_pdf_bytes).decode("utf-8"),
        "filename": "doc-sds-001.pdf",
    }
    reg_resp = client.post("/v1/dossiers/documents/register", json=reg_payload, headers=auth_headers)
    assert reg_resp.status_code == 200

    create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
    assert create_resp.status_code == 200
    create_data = create_resp.json()
    assert create_data["status"] == "created"
    case_id = create_data["case_id"]
    case_digest = create_data["case_digest"]

    # 2. Compile and Run Plan -> Pauses at approval checkpoint with bound approval_request_id
    run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    assert run_resp.status_code == 200
    run_data = run_resp.json()
    plan_digest = run_data["plan_digest"]
    execution = run_data["execution"]
    assert execution["status"] == "awaiting_approval"
    checkpoint = execution["checkpoint"]
    checkpoint_id = checkpoint["checkpoint_id"]
    approval_request_id = execution["approval_request_id"]
    evidence_digests = execution["evidence_digests"]
    assert approval_request_id is not None
    # Verify pending steps include all subsequent steps
    assert "step_human_regulatory_approval" in checkpoint["pending_step_ids"]
    assert "step_assemble_pif_manifest" in checkpoint["pending_step_ids"]

    # 3. Submit Approval Decision with 3 matching digests and bound approval_request_id
    decision_payload = {
        "checkpoint_id": checkpoint_id,
        "run_id": checkpoint["run_id"],
        "approval_request_id": approval_request_id,
        "idempotency_key": "idemp-run-001",
        "decision": "approved",
        "reason": "All toxicology parameters verified against SCCS 12th revision.",
        "case_digest": case_digest,
        "plan_digest": plan_digest,
        "evidence_digests": evidence_digests,
    }

    decide_resp = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert decide_resp.status_code == 200
    decide_data = decide_resp.json()
    assert decide_data["status"] == "decided"
    assert decide_data["decision"] == "approved"
    pdx_resume = decide_data["pdx_resume"]
    assert pdx_resume["status"] == "completed"
    assert pdx_resume["final_manifest"]["status"] == "FINALIZED_COMPLIANT"

    # 4. Check Audit Trail
    run_id = checkpoint["run_id"]
    audit_resp = client.get(f"/v1/audit/runs/{run_id}", headers=auth_headers)
    assert audit_resp.status_code == 200
    audit_events = audit_resp.json()
    event_types = [e["event_type"] for e in audit_events]
    assert "CASE_CREATED" in event_types
    assert "PLAN_COMPILED" in event_types
    assert "CHECKPOINT_CREATED" in event_types

def test_approval_rejection_on_mismatched_approval_request_id_412(client, auth_headers):
    happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
    case_id = create_resp.json()["case_id"]

    run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    execution = run_resp.json()["execution"]
    checkpoint = execution["checkpoint"]

    decision_payload = {
        "checkpoint_id": checkpoint["checkpoint_id"],
        "run_id": checkpoint["run_id"],
        "approval_request_id": str(uuid4()),  # Random mismatched UUID
        "idempotency_key": "idemp-mismatch-req-id",
        "decision": "approved",
        "case_digest": checkpoint["subject_digest"],
        "plan_digest": checkpoint["plan_digest"],
        "evidence_digests": checkpoint["evidence_digests"],
    }

    resp = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp.status_code == 412
    assert "Approval Request ID mismatch" in resp.json()["detail"]

def test_approval_rejection_on_missing_approval_request_412(client, auth_headers):
    from fleet_api.deps import checkpoint_store
    from fleet_governance_core.models.approval import CheckpointStatusEnum, PDXWorkflowCheckpoint
    
    # Save a checkpoint directly without saving any approval request
    chk = PDXWorkflowCheckpoint(
        checkpoint_id="chk-missing-req-001",
        run_id="run-orphaned-001",
        subject_digest="0" * 64,
        plan_digest="1" * 64,
        status=CheckpointStatusEnum.PENDING,
    )
    checkpoint_store.save_checkpoint("tenant-acme-corp", chk)

    decision_payload = {
        "checkpoint_id": "chk-missing-req-001",
        "run_id": "run-orphaned-001",
        "approval_request_id": str(uuid4()),
        "idempotency_key": "idemp-missing-req",
        "decision": "approved",
        "case_digest": "0" * 64,
        "plan_digest": "1" * 64,
        "evidence_digests": {},
    }

    resp = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp.status_code == 412
    assert "Approval request not found for checkpoint" in resp.json()["detail"]


def test_approval_rejection_on_digest_tampering_412(client, auth_headers):
    happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
    case_id = create_resp.json()["case_id"]

    run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    execution = run_resp.json()["execution"]
    checkpoint = execution["checkpoint"]
    approval_request_id = execution["approval_request_id"]

    decision_payload = {
        "checkpoint_id": checkpoint["checkpoint_id"],
        "run_id": checkpoint["run_id"],
        "approval_request_id": approval_request_id,
        "idempotency_key": "idemp-run-tamper",
        "decision": "approved",
        "case_digest": "0" * 64,  # Tampered digest
        "plan_digest": checkpoint["plan_digest"],
        "evidence_digests": checkpoint["evidence_digests"],
    }

    resp = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp.status_code == 412
    assert "Precondition Failed" in resp.json()["detail"]

def test_idempotency_conflict_409(client, auth_headers):
    happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
    case_id = create_resp.json()["case_id"]

    run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    execution = run_resp.json()["execution"]
    checkpoint = execution["checkpoint"]
    approval_request_id = execution["approval_request_id"]

    decision_payload = {
        "checkpoint_id": checkpoint["checkpoint_id"],
        "run_id": checkpoint["run_id"],
        "approval_request_id": approval_request_id,
        "idempotency_key": "idemp-key-shared-01",
        "decision": "approved",
        "case_digest": checkpoint["subject_digest"],
        "plan_digest": checkpoint["plan_digest"],
        "evidence_digests": checkpoint["evidence_digests"],
    }

    # First decision succeeds
    resp1 = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp1.status_code == 200

    # Conflicting decision on same key returns 409
    conflicting_payload = dict(decision_payload)
    conflicting_payload["decision"] = "rejected"
    resp2 = client.post("/v1/approval/decide", json=conflicting_payload, headers=auth_headers)
    assert resp2.status_code == 409
    assert "Idempotency Conflict" in resp2.json()["detail"]

def test_idempotent_replay_with_different_approval_request_id_fails(client, auth_headers):
    happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
    case_id = create_resp.json()["case_id"]

    run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    execution = run_resp.json()["execution"]
    checkpoint = execution["checkpoint"]
    approval_request_id = execution["approval_request_id"]

    decision_payload = {
        "checkpoint_id": checkpoint["checkpoint_id"],
        "run_id": checkpoint["run_id"],
        "approval_request_id": approval_request_id,
        "idempotency_key": f"idemp-replay-req-diff-{case_id[:8]}",
        "decision": "approved",
        "case_digest": checkpoint["subject_digest"],
        "plan_digest": checkpoint["plan_digest"],
        "evidence_digests": checkpoint["evidence_digests"],
    }

    # 1. Initial valid decision succeeds
    resp1 = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp1.status_code == 200

    # 2. Resend with same key and same digests, but different approval_request_id -> MUST fail (409 or 412)
    tampered_payload = dict(decision_payload)
    tampered_payload["approval_request_id"] = str(uuid4())
    resp2 = client.post("/v1/approval/decide", json=tampered_payload, headers=auth_headers)
    assert resp2.status_code in (409, 412)
    assert resp2.status_code != 200

def test_idempotent_replay_with_different_reason_fails(client, auth_headers):
    happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
    case_id = create_resp.json()["case_id"]

    run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    execution = run_resp.json()["execution"]
    checkpoint = execution["checkpoint"]
    approval_request_id = execution["approval_request_id"]

    decision_payload = {
        "checkpoint_id": checkpoint["checkpoint_id"],
        "run_id": checkpoint["run_id"],
        "approval_request_id": approval_request_id,
        "idempotency_key": f"idemp-reason-diff-{case_id[:8]}",
        "decision": "approved",
        "reason": "Original Reason A",
        "case_digest": checkpoint["subject_digest"],
        "plan_digest": checkpoint["plan_digest"],
        "evidence_digests": checkpoint["evidence_digests"],
    }

    # 1. Initial valid decision succeeds
    resp1 = client.post("/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert resp1.status_code == 200

    # 2. Resend with same key and same digests, but modified reason -> MUST fail with 409
    conflicting_payload = dict(decision_payload)
    conflicting_payload["reason"] = "Modified Reason B"
    resp2 = client.post("/v1/approval/decide", json=conflicting_payload, headers=auth_headers)
    assert resp2.status_code == 409
    assert "Idempotency Conflict" in resp2.json()["detail"]

def generate_golden_document_bytes(fmt_ext: str) -> bytes:
    """Generate real binary content parsable by upstream ProDocuX kernel parsers."""
    if fmt_ext == ".pdf":
        return (
            b"%PDF-1.4\n1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
            b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
            b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 100 100] >>\nendobj\n"
            b"xref\n0 4\n0000000000 65535 f \n0000000010 00000 n \n0000000060 00000 n \n0000000117 00000 n \n"
            b"trailer\n<< /Size 4 /Root 1 0 R >>\nstartxref\n193\n%%EOF\n"
        )
    elif fmt_ext == ".docx":
        import io
        from docx import Document
        buf = io.BytesIO()
        doc = Document()
        doc.add_heading("Cosmetics Product Information File", 0)
        doc.add_paragraph("Product: Fortified Night Serum")
        doc.save(buf)
        return buf.getvalue()
    elif fmt_ext == ".csv":
        return b"inci_name,cas_number,percentage\nAqua,7732-18-5,85.0\nGlycerin,56-81-5,5.0\n"
    elif fmt_ext == ".xlsx":
        import io
        from openpyxl import Workbook
        buf = io.BytesIO()
        wb = Workbook()
        ws = wb.active
        ws.title = "Ingredients"
        ws.append(["Component", "Percent"])
        ws.append(["Retinol", 0.05])
        wb.save(buf)
        return buf.getvalue()
    elif fmt_ext == ".pptx":
        import io
        from pptx import Presentation
        buf = io.BytesIO()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(buf)
        return buf.getvalue()
    raise ValueError(f"Unknown format extension: {fmt_ext}")

class RecordingIntakeAdapter(IntakePort):
    """Spy wrapper around real IntakePort that records exact method dispatch and parameters."""
    def __init__(self, target: IntakePort):
        self.target = target
        self.recorded_calls: List[Tuple[str, str, int]] = []

    def extract_pages(self, document_filename: str, document_bytes: bytes, max_pages: int = 50) -> Dict[str, Any]:
        self.recorded_calls.append(("extract_pages", document_filename, len(document_bytes)))
        return self.target.extract_pages(document_filename, document_bytes, max_pages=max_pages)

    def profile_document(self, document_filename: str, document_bytes: bytes) -> Dict[str, Any]:
        self.recorded_calls.append(("profile_document", document_filename, len(document_bytes)))
        return self.target.profile_document(document_filename, document_bytes)

    def profile_table(self, document_filename: str, document_bytes: bytes) -> Dict[str, Any]:
        self.recorded_calls.append(("profile_table", document_filename, len(document_bytes)))
        return self.target.profile_table(document_filename, document_bytes)

    def profile_workbook(self, document_filename: str, document_bytes: bytes) -> Dict[str, Any]:
        self.recorded_calls.append(("profile_workbook", document_filename, len(document_bytes)))
        return self.target.profile_workbook(document_filename, document_bytes)

    def profile_presentation(self, document_filename: str, document_bytes: bytes) -> Dict[str, Any]:
        self.recorded_calls.append(("profile_presentation", document_filename, len(document_bytes)))
        return self.target.profile_presentation(document_filename, document_bytes)

@pytest.mark.parametrize(
    "fmt_ext,doc_id,filename,expected_tool,expected_op",
    [
        (".pdf", "doc-test-pdf", "safety_sheet.pdf", "prodocux.extract_pages", "extract_pages"),
        (".docx", "doc-test-docx", "specification.docx", "prodocux.profile_document", "profile_document"),
        (".csv", "doc-test-csv", "table_data.csv", "prodocux.profile_table", "profile_table"),
        (".xlsx", "doc-test-xlsx", "toxicology.xlsx", "prodocux.profile_workbook", "profile_workbook"),
        (".pptx", "doc-test-pptx", "presentation.pptx", "prodocux.profile_presentation", "profile_presentation"),
    ],
)
def test_five_formats_api_compile_and_run_lifecycle(client, auth_headers, fmt_ext, doc_id, filename, expected_tool, expected_op):
    """Verify that all five supported document formats complete the full API compile, register, and run cycle through LivePDXCoreOrchestrator and in-process ProDocuX parsers."""
    import sys
    from pathlib import Path
    from starlette.testclient import TestClient
    from fleet_adapter_prodocux import ProDocuXHttpIntakeAdapter
    from fleet_adapter_pdx.orchestrator import LivePDXCoreOrchestrator
    from fleet_api import deps
    from fleet_api.main import app

    import os
    try:
        from api.main import app as prodocux_app
    except ImportError:
        prodocux_repo_env = os.getenv("PRODOCUX_REPO_DIR")
        if prodocux_repo_env and Path(prodocux_repo_env).exists():
            if prodocux_repo_env not in sys.path:
                sys.path.insert(0, prodocux_repo_env)
            try:
                from api.main import app as prodocux_app
            except ImportError:
                prodocux_app = None
        else:
            prodocux_app = None

    if prodocux_app is None:
        pytest.skip("ProDocuX package or PRODOCUX_REPO_DIR not installed/available (NOT RUN)")

    # Wire LivePDXCoreOrchestrator with recording spy wrapping real in-process ProDocuX TestClient kernel
    real_inprocess_intake = ProDocuXHttpIntakeAdapter(
        base_url="https://prodocux.internal",
        http_client=TestClient(prodocux_app),
    )
    spy_intake = RecordingIntakeAdapter(real_inprocess_intake)
    live_orch = LivePDXCoreOrchestrator(
        approval_ledger=deps.shared_approval_ledger,
        intake_adapter=spy_intake,
        document_resolver=deps.document_resolver,
    )
    
    # Inject live orchestrator via FastAPI dependency overrides
    app.dependency_overrides[deps.get_orchestrator] = lambda: live_orch

    try:
        # 1. Generate real golden binary content and register
        golden_bytes = generate_golden_document_bytes(fmt_ext)
        sample_sha256 = hashlib.sha256(golden_bytes).hexdigest()

        reg_payload = {
            "doc_id": doc_id,
            "content_b64": base64.b64encode(golden_bytes).decode("utf-8"),
            "filename": filename,
        }
        reg_resp = client.post("/v1/dossiers/documents/register", json=reg_payload, headers=auth_headers)
        assert reg_resp.status_code == 200
        assert reg_resp.json()["status"] == "registered"

        # 2. Create Dossier Case with authoritative document metadata
        happy_raw = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
        happy_raw["case_id"] = str(uuid4())
        happy_raw["supplier_documents"] = [
            {
                "doc_id": doc_id,
                "filename": filename,
                "doc_type": "SDS",
                "sha256": sample_sha256,
                "supplier_name": "Format Test Supplier",
                "issue_date": "2025-01-10",
                "expiry_date": "2028-01-10",
            }
        ]

        create_resp = client.post("/v1/dossiers/create", json=happy_raw, headers=auth_headers)
        assert create_resp.status_code == 200
        case_id = create_resp.json()["case_id"]

        # 3. Compile and Run Dossier -> Executes LivePDXCoreOrchestrator against real in-process ProDocuX parser
        run_resp = client.post(f"/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
        assert run_resp.status_code == 200
        run_data = run_resp.json()
        assert run_data["execution"]["status"] == "awaiting_approval"
        
        # Verify plan step tool mapping
        plan_steps = run_data["plan"]["steps"]
        extract_step = [s for s in plan_steps if s.get("inputs", {}).get("document_id") == doc_id][0]
        assert extract_step["tool"] == expected_tool
        assert extract_step["inputs"]["document_filename"] == filename

        # Verify real intake evidence digest was generated
        evidence_digests = run_data["execution"]["evidence_digests"]
        assert f"{extract_step['id']}_output.json" in evidence_digests

        # Explicitly assert that the real ProDocuX parser method was dispatched with exact filename and golden byte count
        assert len(spy_intake.recorded_calls) == 1
        recorded_op, recorded_filename, recorded_len = spy_intake.recorded_calls[0]
        assert recorded_op == expected_op
        assert recorded_filename == filename
        assert recorded_len == len(golden_bytes)
    finally:
        app.dependency_overrides.pop(deps.get_orchestrator, None)


def test_evidence_package_endpoint(client, auth_headers):
    """Verify that GET /v1/evidence/runs/{run_id} returns checksummed evidence package with canonical SHA-256."""
    from fleet_api.deps import audit_log
    from fleet_governance_core.models.audit import AuditEvent, AuditEventTypeEnum

    audit_log.append_audit_event(
        AuditEvent(
            tenant_id="tenant-acme-corp",
            run_id="run-test-12345",
            event_type=AuditEventTypeEnum.PLAN_COMPILED,
            actor_id="usr-safety-officer-01",
            payload={"plan_digest": "a" * 64, "case_digest": "b" * 64},
        )
    )

    resp = client.get("/v1/evidence/runs/run-test-12345", headers=auth_headers)
    assert resp.status_code == 200
    data = resp.json()
    assert data["package_type"] == "checksummed_evidence_package"
    assert data["version"] == "0.3.2"
    assert data["tenant_id"] == "tenant-acme-corp"
    assert data["run_id"] == "run-test-12345"
    assert data["integrity"] == "sha256_checksum_only"
    assert data["digitally_signed"] is False
    assert data["artifact_store_mode"] == "local_filesystem_ephemeral"
    assert data["case_digest"] == "b" * 64
    assert data["plan_digest"] == "a" * 64
    assert "package_sha256" in data
    assert len(data["package_sha256"]) == 64

    # Fail closed on non-existent run
    resp_404 = client.get("/v1/evidence/runs/unknown-run-999", headers=auth_headers)
    assert resp_404.status_code == 404






