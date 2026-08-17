"""
Gate B9-Docker-Live: Docker Production Live-Adapter Deployment Gate & Full End-to-End Lifecycle Conformance Suite.
Validates:
1. Pinned Git HEAD Docker Image execution in strict FLEET_ENV=production runtime.
2. Containerized Live ProDocuX HTTP Server deployment on isolated Docker bridge network.
3. G6A Live HTTP Endpoint Conformance (version, capabilities, and 5-format document extraction over HTTP wire).
4. Live Intake Adapter (ProDocuXHttpIntakeAdapter) + Live PDX Core Orchestrator (LivePDXCoreOrchestrator).
5. Production Liveness (/v1/health) and Readiness (/v1/ready) Probes under live adapters.
6. Cryptographic Multi-Tenant Isolation in production environment.
7. Full 5-Format Live Lifecycle: PDF, DOCX, CSV, XLSX, PPTX registration, compilation, checkpointing, and single-transaction CSO approval.
8. Container Hard Termination (docker kill) & Persistent Volume Restart Recovery (SQLite ACID + Local Artifact Store).
9. Idempotent Approval Replay Verification on restarted production container.
"""
import base64
import hashlib
import io
import json
import os
from pathlib import Path
import socket
import subprocess
import sys
import time
from typing import Dict, Optional, Tuple
from uuid import uuid4

import jwt
from openpyxl import Workbook
from pypdf import PdfWriter
import pytest
import requests

ROOT_DIR = Path(__file__).resolve().parent.parent
FIXTURES_DIR = ROOT_DIR / "fixtures"

JWT_SECRET = "b9-docker-live-prod-secret-2026-fortified-998877665544332211"

EXPECTED_PDX_COMMIT = "55a9293c8d5c0091e04e457dc43f662058e50068"
EXPECTED_PRODOCUX_COMMIT = "c8acd2ba69c23458cb2589d8450246fe9b16424f"
EXPECTED_MANIFEST_SHA = "a5eff2cc21aeff8eb0f6cad1e6e7dd3f50daff3ea3faedb4989c03b1af87161c"


def get_current_git_commit() -> str:
    res = subprocess.run(["git", "rev-parse", "HEAD"], cwd=str(ROOT_DIR), capture_output=True, text=True, check=True)
    return res.stdout.strip()


def get_free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


def make_jwt_token(tenant_id: str, sub: str, email: str, role: str) -> str:
    now = int(time.time())
    payload = {
        "iss": "fortified-enterprise-fleet-auth",
        "sub": sub,
        "tenant_id": tenant_id,
        "roles": [role],
        "email": email,
        "iat": now,
        "exp": now + 3600,
    }
    return jwt.encode(payload, JWT_SECRET, algorithm="HS256")


def wait_for_server(url: str, timeout_seconds: float = 20.0) -> bool:
    start = time.time()
    while time.time() - start < timeout_seconds:
        try:
            resp = requests.get(url, timeout=1.0)
            if resp.status_code == 200:
                return True
        except Exception:
            pass
        time.sleep(0.3)
    return False


@pytest.fixture(scope="session")
def docker_live_environment():
    """
    Session fixture that:
    1. Builds the OCI revision-pinned Docker image.
    2. Creates an isolated Docker bridge network.
    3. Starts the containerized ProDocuX HTTP Server.
    """
    head_commit = get_current_git_commit()
    image_tag = f"fortifiedreg-fleet:live-test-{head_commit[:8]}"
    network_name = f"fleet-live-net-{uuid4().hex[:6]}"
    prodocux_container = f"prodocux-live-{uuid4().hex[:6]}"
    prodocux_host_port = get_free_port()

    # 1. Build test image
    build_cmd = [
        "docker",
        "build",
        "--build-arg",
        f"GIT_COMMIT={head_commit}",
        "-t",
        image_tag,
        str(ROOT_DIR),
    ]
    subprocess.run(build_cmd, check=True, capture_output=True, text=True)

    # 2. Create Docker bridge network
    subprocess.run(["docker", "network", "create", network_name], check=True, capture_output=True, text=True)

    # 3. Start ProDocuX HTTP Server container on network
    start_prodocux_cmd = [
        "docker",
        "run",
        "-d",
        "--name",
        prodocux_container,
        "--network",
        network_name,
        "-p",
        f"{prodocux_host_port}:8900",
        "--entrypoint",
        "uvicorn",
        image_tag,
        "api.main:app",
        "--host",
        "0.0.0.0",
        "--port",
        "8900",
    ]
    subprocess.run(start_prodocux_cmd, check=True, capture_output=True, text=True)

    prodocux_host_url = f"http://127.0.0.1:{prodocux_host_port}"
    prodocux_container_url = f"http://{prodocux_container}:8900"

    if not wait_for_server(f"{prodocux_host_url}/v1/version", timeout_seconds=15.0):
        logs = subprocess.run(["docker", "logs", prodocux_container], capture_output=True, text=True).stdout
        raise RuntimeError(f"ProDocuX live container failed to start. Logs:\n{logs}")

    yield {
        "image_tag": image_tag,
        "head_commit": head_commit,
        "network_name": network_name,
        "prodocux_container": prodocux_container,
        "prodocux_host_url": prodocux_host_url,
        "prodocux_container_url": prodocux_container_url,
    }

    # Teardown
    subprocess.run(["docker", "rm", "-f", prodocux_container], capture_output=True, text=True)
    subprocess.run(["docker", "network", "rm", network_name], capture_output=True, text=True)
    subprocess.run(["docker", "rmi", "-f", image_tag], capture_output=True, text=True)


class FleetProductionContainerProcess:
    def __init__(
        self,
        image_tag: str,
        network_name: str,
        prodocux_url: str,
        data_dir: Path,
        artifacts_dir: Path,
        port: int,
    ):
        self.image_tag = image_tag
        self.network_name = network_name
        self.prodocux_url = prodocux_url
        self.data_dir = data_dir
        self.artifacts_dir = artifacts_dir
        self.port = port
        self.base_url = f"http://127.0.0.1:{port}"
        self.container_id = None
        self.container_name = f"fleet-prod-test-{uuid4().hex[:8]}"

    def start(self):
        cmd = [
            "docker",
            "run",
            "-d",
            "--name",
            self.container_name,
            "--network",
            self.network_name,
            "-p",
            f"{self.port}:8000",
            "-e",
            f"FLEET_JWT_SECRET={JWT_SECRET}",
            "-e",
            "FLEET_ENV=production",
            "-e",
            "FLEET_INTAKE_ADAPTER=live",
            "-e",
            "FLEET_PDX_ADAPTER=live",
            "-e",
            f"PRODOCUX_BASE_URL={self.prodocux_url}",
            "-v",
            f"{self.data_dir.resolve()}:/app/data",
            "-v",
            f"{self.artifacts_dir.resolve()}:/app/artifacts",
            self.image_tag,
        ]
        res = subprocess.run(cmd, capture_output=True, text=True, check=True)
        self.container_id = res.stdout.strip()

        if not wait_for_server(f"{self.base_url}/v1/health", timeout_seconds=20.0):
            logs = self.read_logs()
            self.stop()
            raise RuntimeError(f"Fleet Production container failed to start within timeout. Logs:\n{logs}")

    def kill(self):
        if self.container_name:
            subprocess.run(["docker", "kill", self.container_name], capture_output=True, text=True)
            subprocess.run(["docker", "rm", "-f", self.container_name], capture_output=True, text=True)
            self.container_id = None

    def stop(self):
        if self.container_name:
            subprocess.run(["docker", "rm", "-f", self.container_name], capture_output=True, text=True)
            self.container_id = None

    def read_logs(self) -> str:
        if self.container_name:
            res = subprocess.run(["docker", "logs", self.container_name], capture_output=True, text=True)
            return res.stdout + res.stderr
        return ""


@pytest.fixture
def fleet_production_env(docker_live_environment, tmp_path):
    env_info = docker_live_environment
    port = get_free_port()
    data_dir = tmp_path / "prod_data"
    data_dir.mkdir(parents=True, exist_ok=True)
    artifacts_dir = tmp_path / "prod_artifacts"
    artifacts_dir.mkdir(parents=True, exist_ok=True)

    container = FleetProductionContainerProcess(
        image_tag=env_info["image_tag"],
        network_name=env_info["network_name"],
        prodocux_url=env_info["prodocux_container_url"],
        data_dir=data_dir,
        artifacts_dir=artifacts_dir,
        port=port,
    )
    container.start()
    try:
        yield container, env_info
    finally:
        container.stop()


def generate_valid_document_bytes() -> Dict[str, Tuple[str, str, bytes]]:
    """Generate 5 valid binary documents for registration."""
    # 1. Valid PDF
    pdf_buf = io.BytesIO()
    pdf_writer = PdfWriter()
    pdf_writer.add_blank_page(width=100, height=100)
    pdf_writer.write(pdf_buf)
    pdf_bytes = pdf_buf.getvalue()

    # 2. Valid DOCX
    from docx import Document
    doc = Document()
    doc.add_heading("Specification", 0)
    doc.add_paragraph("Valid docx specification content.")
    buf_docx = io.BytesIO()
    doc.save(buf_docx)
    docx_bytes = buf_docx.getvalue()

    # 3. Valid CSV
    csv_bytes = b"inci_name,cas_number,percentage\nAqua,7732-18-5,85.0\nGlycerin,56-81-5,5.0\n"

    # 4. Valid XLSX
    xlsx_buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "Ingredients"
    ws.append(["Component", "Percent"])
    ws.append(["Retinol", 0.05])
    wb.save(xlsx_buf)
    xlsx_bytes = xlsx_buf.getvalue()

    # 5. Valid PPTX
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentation Title"
    buf_pptx = io.BytesIO()
    prs.save(buf_pptx)
    pptx_bytes = buf_pptx.getvalue()

    return {
        ".pdf": ("doc-sds-001", "safety_sheet.pdf", pdf_bytes),
        ".docx": ("doc-spec-001", "specification.docx", docx_bytes),
        ".csv": ("doc-table-001", "formulation.csv", csv_bytes),
        ".xlsx": ("doc-tox-001", "toxicology.xlsx", xlsx_bytes),
        ".pptx": ("doc-pres-001", "presentation.pptx", pptx_bytes),
    }


# ---------------------------------------------------------------------------
# B9 Docker Production Live-Adapter Gate Tests
# ---------------------------------------------------------------------------

def test_b9_g6a_live_http_conformance_against_container(docker_live_environment):
    """Gate G6A: Direct HTTP conformance against containerized live ProDocuX server."""
    env_info = docker_live_environment
    host_url = env_info["prodocux_host_url"]

    # 1. Version Check
    ver_resp = requests.get(f"{host_url}/v1/version")
    assert ver_resp.status_code == 200
    ver = ver_resp.json()
    assert ver.get("kernel_version") == "0.2.0"
    assert ver.get("api_version") == "v1"

    # 2. Capabilities Check
    cap_resp = requests.get(f"{host_url}/v1/intake/capabilities")
    assert cap_resp.status_code == 200
    cap = cap_resp.json()
    assert cap.get("schema_version") == "prodocux_intake_capabilities_v1"

    # 3. Live 5-Format Ingestion over HTTP wire
    docs = generate_valid_document_bytes()

    # PDF
    _, _, pdf_bytes = docs[".pdf"]
    pdf_req = {"document_filename": "live.pdf", "document_b64": base64.b64encode(pdf_bytes).decode(), "max_pages": 50}
    pdf_resp = requests.post(f"{host_url}/v1/intake/extract-pages", json=pdf_req)
    assert pdf_resp.status_code == 200
    assert pdf_resp.json().get("status") in ("success", "ocr_required")

    # DOCX
    _, _, docx_bytes = docs[".docx"]
    docx_req = {"document_filename": "live.docx", "document_b64": base64.b64encode(docx_bytes).decode()}
    docx_resp = requests.post(f"{host_url}/v1/intake/profile-document", json=docx_req)
    assert docx_resp.status_code == 200
    assert docx_resp.json().get("profile", {}).get("schema_version") == "prodocux_docx_profile_v1"

    # CSV
    _, _, csv_bytes = docs[".csv"]
    csv_req = {"document_filename": "live.csv", "document_b64": base64.b64encode(csv_bytes).decode()}
    csv_resp = requests.post(f"{host_url}/v1/intake/profile-table", json=csv_req)
    assert csv_resp.status_code == 200
    assert csv_resp.json().get("profile", {}).get("schema_version") == "prodocux_table_profile_v1"

    # XLSX
    _, _, xlsx_bytes = docs[".xlsx"]
    xlsx_req = {"document_filename": "live.xlsx", "document_b64": base64.b64encode(xlsx_bytes).decode()}
    xlsx_resp = requests.post(f"{host_url}/v1/intake/profile-workbook", json=xlsx_req)
    assert xlsx_resp.status_code == 200
    assert xlsx_resp.json().get("profile", {}).get("schema_version") == "prodocux_workbook_profile_v1"

    # PPTX
    _, _, pptx_bytes = docs[".pptx"]
    pptx_req = {"document_filename": "live.pptx", "document_b64": base64.b64encode(pptx_bytes).decode()}
    pptx_resp = requests.post(f"{host_url}/v1/intake/profile-presentation", json=pptx_req)
    assert pptx_resp.status_code == 200
    assert pptx_resp.json().get("profile", {}).get("schema_version") == "prodocux_presentation_profile_v1"


def test_b9_docker_production_probes_and_live_adapters(fleet_production_env):
    """Verify health & ready probes in FLEET_ENV=production with live adapters."""
    container, _ = fleet_production_env

    # 1. Health Probe
    resp_health = requests.get(f"{container.base_url}/v1/health")
    assert resp_health.status_code == 200
    assert resp_health.json()["status"] == "healthy"
    assert resp_health.json()["environment"] == "production"

    # 2. Ready Probe
    resp_ready = requests.get(f"{container.base_url}/v1/ready")
    assert resp_ready.status_code == 200
    data_ready = resp_ready.json()
    assert data_ready["status"] == "ready"
    assert data_ready["adapters"]["intake_mode"] == "live"
    assert data_ready["adapters"]["pdx_mode"] == "live"


def test_b9_docker_production_multi_tenant_isolation(fleet_production_env):
    """Verify cryptographic multi-tenant isolation under production configuration."""
    container, _ = fleet_production_env
    token_tenant_a = make_jwt_token("tenant-acme-corp", "usr-cso-a", "cso@acme.com", "cso")
    token_tenant_b = make_jwt_token("tenant-globex-inc", "usr-cso-b", "cso@globex.com", "cso")

    headers_a = {"Authorization": f"Bearer {token_tenant_a}"}
    headers_b = {"Authorization": f"Bearer {token_tenant_b}"}

    case_a_id = str(uuid4())
    raw_case_a = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    raw_case_a["case_id"] = case_a_id
    raw_case_a["tenant_id"] = "tenant-acme-corp"

    create_resp = requests.post(f"{container.base_url}/v1/dossiers/create", json=raw_case_a, headers=headers_a)
    assert create_resp.status_code == 200

    # Cross-tenant 404
    get_by_b = requests.get(f"{container.base_url}/v1/dossiers/{case_a_id}", headers=headers_b)
    assert get_by_b.status_code == 404

    # Cross-tenant 403
    bad_case = dict(raw_case_a, case_id=str(uuid4()), tenant_id="tenant-acme-corp")
    bad_create = requests.post(f"{container.base_url}/v1/dossiers/create", json=bad_case, headers=headers_b)
    assert bad_create.status_code == 403


def test_b9_docker_production_live_five_formats_lifecycle_and_volume_restart(fleet_production_env):
    """
    Full End-to-End Production Live-Adapter Lifecycle:
    1. Register 5 supplier document formats.
    2. Compile & run with Live PDX Core + Live ProDocuX HTTP Intake over network.
    3. Verify live extraction, verifier bridges, and checkpoint generation.
    4. Submit CSO approval decision -> publish final artifact.
    5. Hard kill container (docker kill).
    6. Restart on persistent volume -> verify idempotent replay and state retrieval.
    """
    container, env_info = fleet_production_env
    token_cso = make_jwt_token("tenant-acme-corp", "usr-cso-1", "cso@acme.com", "cso")
    auth_headers = {"Authorization": f"Bearer {token_cso}"}

    # 1. Register 5 document formats
    docs_map = generate_valid_document_bytes()
    doc_types = ["SDS", "COA", "GMP_CERT", "IFRA_CERT", "COA"]
    registered_docs = []

    for i, (ext, (doc_id, filename, raw_bytes)) in enumerate(docs_map.items()):
        reg_resp = requests.post(
            f"{container.base_url}/v1/dossiers/documents/register",
            json={
                "doc_id": doc_id,
                "content_b64": base64.b64encode(raw_bytes).decode(),
                "filename": filename,
            },
            headers=auth_headers,
        )
        assert reg_resp.status_code == 200, f"Registration failed for {filename}: {reg_resp.text}"
        registered_docs.append({
            "doc_id": doc_id,
            "filename": filename,
            "doc_type": doc_types[i],
            "sha256": hashlib.sha256(raw_bytes).hexdigest(),
            "supplier_name": "BioSynthetics Ltd",
            "issue_date": "2025-01-10",
            "expiry_date": "2028-01-10",
        })

    # 2. Create Dossier Case
    case_id = str(uuid4())
    raw_case = json.loads((FIXTURES_DIR / "c2_dossier_case_happy_path.json").read_text(encoding="utf-8"))["data"]
    raw_case["case_id"] = case_id
    raw_case["tenant_id"] = "tenant-acme-corp"
    raw_case["supplier_documents"] = registered_docs

    create_resp = requests.post(f"{container.base_url}/v1/dossiers/create", json=raw_case, headers=auth_headers)
    assert create_resp.status_code == 200

    # 3. Compile and Run with LIVE PDX Core + LIVE ProDocuX HTTP Intake
    run_resp = requests.post(f"{container.base_url}/v1/dossiers/{case_id}/compile-and-run", headers=auth_headers)
    assert run_resp.status_code == 200, f"Compile and run failed: {run_resp.text}"
    exec_data = run_resp.json()["execution"]
    assert exec_data["status"] == "awaiting_approval"

    checkpoint = exec_data["checkpoint"]
    approval_request_id = exec_data["approval_request_id"]

    # 4. Submit CSO Approval Decision
    idempotency_key = f"prod-live-idemp-{case_id[:8]}"
    decision_payload = {
        "checkpoint_id": checkpoint["checkpoint_id"],
        "run_id": checkpoint["run_id"],
        "approval_request_id": approval_request_id,
        "idempotency_key": idempotency_key,
        "decision": "approved",
        "reason": "Production live-adapter deployment authorization",
        "case_digest": checkpoint["subject_digest"],
        "plan_digest": checkpoint["plan_digest"],
        "evidence_digests": checkpoint["evidence_digests"],
    }

    dec_resp = requests.post(f"{container.base_url}/v1/approval/decide", json=decision_payload, headers=auth_headers)
    assert dec_resp.status_code == 200, f"Approval decision failed: {dec_resp.text}"
    dec_data = dec_resp.json()
    assert dec_data["status"] == "decided"
    assert dec_data["decision"] == "approved"
    artifact_ident = dec_data["artifact_identity"]
    storage_uri = artifact_ident.get("uri") or artifact_ident.get("storage_uri")
    assert storage_uri and storage_uri.startswith("artifact://")

    # 5. Hard kill container
    data_dir = container.data_dir
    artifacts_dir = container.artifacts_dir
    port = container.port
    container.kill()

    # 6. Start new container on same persistent volume
    new_container = FleetProductionContainerProcess(
        image_tag=env_info["image_tag"],
        network_name=env_info["network_name"],
        prodocux_url=env_info["prodocux_container_url"],
        data_dir=data_dir,
        artifacts_dir=artifacts_dir,
        port=port,
    )
    new_container.start()

    try:
        # Re-check case retrieval
        get_resp = requests.get(f"{new_container.base_url}/v1/dossiers/{case_id}", headers=auth_headers)
        assert get_resp.status_code == 200
        assert get_resp.json()["case"]["case_id"] == case_id

        # Idempotent replay
        replay_resp = requests.post(f"{new_container.base_url}/v1/approval/decide", json=decision_payload, headers=auth_headers)
        assert replay_resp.status_code == 200
        assert replay_resp.json()["status"] == "decided"
        assert replay_resp.json()["is_idempotent_replay"] is True
        assert replay_resp.json()["artifact_identity"]["sha256"] == artifact_ident["sha256"]

        # Conflict rejection on altered payload
        conflict_payload = dict(decision_payload, reason="Altered Reason After Restart")
        conflict_resp = requests.post(f"{new_container.base_url}/v1/approval/decide", json=conflict_payload, headers=auth_headers)
        assert conflict_resp.status_code == 409

    finally:
        new_container.stop()
