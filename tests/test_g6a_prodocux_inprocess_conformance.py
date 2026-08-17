"""
Gate 6A: ProDocuX HTTP Intake In-Process Conformance Tests.
Validates:
1. Adapter pre-flight byte limit boundaries (MAX vs MAX + 1) for all 5 supported formats without network calls.
2. URL validation rules (strict HTTPS in production, forbidden schemes, userinfo, query, fragments).
3. In-process parser integration with ProDocuX FastAPI Kernel via injected TestClient across PDF, DOCX, CSV, XLSX, PPTX.
"""
import io
import sys
from pathlib import Path
import pytest
from starlette.testclient import TestClient

# Ensure upstream ProDocuX repository is accessible for in-process TestClient binding
PRODOCUX_REPO = Path("D:/ProDocuX/prodocux")
if str(PRODOCUX_REPO) not in sys.path:
    sys.path.insert(0, str(PRODOCUX_REPO))

from api.main import app as prodocux_app
from fleet_adapter_prodocux import (
    FORMAT_LIMITS,
    MAX_DOCX_BYTES,
    MAX_PDF_BYTES,
    MAX_PRESENTATION_BYTES,
    MAX_TABLE_BYTES,
    MAX_WORKBOOK_BYTES,
    IntakeConfigurationError,
    IntakePayloadError,
    ProDocuXHttpIntakeAdapter,
    validate_prodocux_url,
)

# Mock/Stub client to intercept calls during size boundary tests
class RecordingMockClient:
    def __init__(self):
        self.calls = []

    def post(self, url, json=None, timeout=None):
        self.calls.append({"url": url, "json": json})
        class DummyResp:
            status_code = 200
            def json(self):
                return {"status": "success", "mock": True}
        return DummyResp()

    def get(self, url, timeout=None):
        self.calls.append({"url": url})
        class DummyResp:
            status_code = 200
            def json(self):
                if "capabilities" in url:
                    return {
                        "schema_version": "prodocux_intake_capabilities_v1",
                        "kernel_version": "0.2.0",
                        "formats": [{"extensions": [".pdf"]}],
                    }
                return {"kernel_version": "0.2.0", "api_version": "v1"}
        return DummyResp()

# ---------------------------------------------------------------------------
# Tier 1: Adapter Size Boundary Unit Tests (MAX vs MAX + 1)
# ---------------------------------------------------------------------------

@pytest.mark.parametrize(
    "filename,max_bytes,method_name",
    [
        ("test.pdf", MAX_PDF_BYTES, "extract_pages"),
        ("test.docx", MAX_DOCX_BYTES, "profile_document"),
        ("test.csv", MAX_TABLE_BYTES, "profile_table"),
        ("test.xlsx", MAX_WORKBOOK_BYTES, "profile_workbook"),
        ("test.pptx", MAX_PRESENTATION_BYTES, "profile_presentation"),
    ],
)
def test_adapter_exact_size_boundary_passes_at_max(filename, max_bytes, method_name):
    mock_client = RecordingMockClient()
    adapter = ProDocuXHttpIntakeAdapter(base_url="http://testserver", http_client=mock_client)
    
    # Payload of exact MAX bytes
    payload = b"0" * max_bytes
    caller = getattr(adapter, method_name)
    res = caller(filename, payload)

    assert res["status"] == "success"
    assert len(mock_client.calls) == 1

@pytest.mark.parametrize(
    "filename,max_bytes,method_name",
    [
        ("test.pdf", MAX_PDF_BYTES, "extract_pages"),
        ("test.docx", MAX_DOCX_BYTES, "profile_document"),
        ("test.csv", MAX_TABLE_BYTES, "profile_table"),
        ("test.xlsx", MAX_WORKBOOK_BYTES, "profile_workbook"),
        ("test.pptx", MAX_PRESENTATION_BYTES, "profile_presentation"),
    ],
)
def test_adapter_size_boundary_rejected_at_max_plus_one(filename, max_bytes, method_name):
    mock_client = RecordingMockClient()
    adapter = ProDocuXHttpIntakeAdapter(base_url="http://testserver", http_client=mock_client)
    
    # Payload of MAX + 1 bytes -> MUST reject before making any network call
    payload = b"0" * (max_bytes + 1)
    caller = getattr(adapter, method_name)
    
    with pytest.raises(IntakePayloadError, match="exceeds effective limit"):
        caller(filename, payload)

    assert len(mock_client.calls) == 0  # Pre-flight check prevents network transmission

def test_empty_payload_rejected():
    mock_client = RecordingMockClient()
    adapter = ProDocuXHttpIntakeAdapter(base_url="http://testserver", http_client=mock_client)
    with pytest.raises(IntakePayloadError, match="empty"):
        adapter.extract_pages("test.pdf", b"")
    assert len(mock_client.calls) == 0

def test_unsupported_extension_rejected():
    mock_client = RecordingMockClient()
    adapter = ProDocuXHttpIntakeAdapter(base_url="http://testserver", http_client=mock_client)
    with pytest.raises(IntakePayloadError, match="Unsupported document format"):
        adapter.extract_pages("test.exe", b"dummy content")
    assert len(mock_client.calls) == 0

# ---------------------------------------------------------------------------
# Tier 1.5: URL Validation Unit Tests
# ---------------------------------------------------------------------------

def test_url_validation_success():
    assert validate_prodocux_url("http://localhost:8900") == "http://localhost:8900"
    assert validate_prodocux_url("https://prodocux.acme.internal:8900/") == "https://prodocux.acme.internal:8900"

def test_url_validation_forbidden_schemes():
    with pytest.raises(IntakeConfigurationError, match="Invalid URL scheme"):
        validate_prodocux_url("file:///etc/passwd")

def test_url_validation_forbidden_userinfo():
    with pytest.raises(IntakeConfigurationError, match="userinfo are forbidden"):
        validate_prodocux_url("http://admin:secret@localhost:8900")

def test_url_validation_forbidden_query_and_fragments():
    with pytest.raises(IntakeConfigurationError, match="Query parameters"):
        validate_prodocux_url("http://localhost:8900?token=xyz")
    with pytest.raises(IntakeConfigurationError, match="Query parameters"):
        validate_prodocux_url("http://localhost:8900#section")

def test_url_validation_production_https_requirement():
    with pytest.raises(IntakeConfigurationError, match="HTTPS is strictly required"):
        validate_prodocux_url("http://prodocux.acme.internal:8900", is_production=True)

# ---------------------------------------------------------------------------
# Tier 2: In-Process Parser Integration Tests with Real ProDocuX Kernel
# ---------------------------------------------------------------------------

@pytest.fixture
def inprocess_adapter():
    client = TestClient(prodocux_app)
    return ProDocuXHttpIntakeAdapter(base_url="http://testserver", http_client=client, is_production=False)

def test_inprocess_version_and_capabilities(inprocess_adapter):
    ver = inprocess_adapter.get_version()
    assert ver["kernel_version"] == "0.2.0"
    assert ver["api_version"] == "v1"

    ready = inprocess_adapter.check_readiness()
    assert ready["status"] == "ready"
    assert ready["schema_version"] == "prodocux_intake_capabilities_v1"
    assert len(ready["formats"]) >= 4


def test_inprocess_valid_pdf_intake(inprocess_adapter):
    from pypdf import PdfWriter
    buf = io.BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    writer.write(buf)
    pdf_bytes = buf.getvalue()

    resp = inprocess_adapter.extract_pages("safety_assessment.pdf", pdf_bytes, max_pages=10)
    assert resp["status"] in ("success", "ocr_required")
    assert resp["page_count"] == 1
    assert "source_sha256" in resp
    assert len(resp["pages"]) == 1

def test_inprocess_valid_docx_intake(inprocess_adapter):
    from docx import Document
    buf = io.BytesIO()
    doc = Document()
    doc.add_heading("Cosmetics Product Information File", 0)
    doc.add_paragraph("Product: Fortified Night Serum")
    doc.save(buf)
    docx_bytes = buf.getvalue()

    resp = inprocess_adapter.profile_document("pif_spec.docx", docx_bytes)
    assert "profile" in resp
    profile = resp["profile"]
    assert profile["schema_version"] == "prodocux_docx_profile_v1"
    assert profile["source"]["name"] == "pif_spec.docx"
    assert profile["paragraph_count"] >= 1

def test_inprocess_valid_csv_intake(inprocess_adapter):
    csv_bytes = b"inci_name,cas_number,percentage\nAqua,7732-18-5,85.0\nGlycerin,56-81-5,5.0\n"
    resp = inprocess_adapter.profile_table("formulation.csv", csv_bytes)
    assert "profile" in resp
    profile = resp["profile"]
    assert profile["row_count"] >= 2
    assert "inci_name" in profile["columns"]

def test_inprocess_valid_xlsx_intake(inprocess_adapter):
    from openpyxl import Workbook
    buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "Ingredients"
    ws.append(["Component", "Percent"])
    ws.append(["Retinol", 0.05])
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    resp = inprocess_adapter.profile_workbook("raw_materials.xlsx", xlsx_bytes)
    assert "profile" in resp
    profile = resp["profile"]
    sheet_names = [s["name"] for s in profile["sheets"]]
    assert "Ingredients" in sheet_names


def test_inprocess_valid_pptx_intake(inprocess_adapter):
    from pptx import Presentation
    buf = io.BytesIO()
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    resp = inprocess_adapter.profile_presentation("regulatory_briefing.pptx", pptx_bytes)
    assert "profile" in resp
    profile = resp["profile"]
    assert profile["slide_count"] == 1
